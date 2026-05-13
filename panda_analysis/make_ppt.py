# -*- coding: utf-8 -*-
"""
make_ppt.py — 课程汇报 PPT (editorial / minimalist 风格)
========================================================
设计语言:
    · 纸白底 + 油墨黑 + 单一点睛色(深红 / 赭石)
    · 编辑级排版:大字号 hero、留白、偏置、发丝线
    · 大汉字水印作为页面节奏锚点
    · 微号 letter-spaced 章节编号 / 页码
共 9 页,16:9。

产出: {学号}_{姓名}_演示PPT.pptx
"""

import os
import sys
from pathlib import Path

import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

sys.stdout.reconfigure(encoding="utf-8")
sys.stderr.reconfigure(encoding="utf-8")


# ============================================================
# 学生信息
# ============================================================
STUDENT_ID    = "202422480546"
STUDENT_NAME  = "罗靖坤"
STUDENT_MAJOR = "旅游管理"
STUDENT_COLLEGE = "旅游学院"
TEACHER_NAME  = "徐 欣"
SUBMIT_DATE   = "2026 / 05"

OUT_DIR = Path("output")
PPT_PATH = Path(f"{STUDENT_ID}_{STUDENT_NAME}_演示PPT.pptx")
STATS = {}


def _clean_dim_name(name):
    return str(name).split(".", 1)[-1].strip() if "." in str(name) else str(name)


def _signed_score(value):
    return f"{float(value):+.2f}".replace("-", "−")


def _load_stats():
    stats = {
        "sample_n": 2000,
        "n_dom": 1000,
        "n_intl": 1000,
        "date_range": "2024 – 2026",
        "pct_translated_intl": 0.0,
        "pct_translated": 0.0,
        "total_tokens": 0,
        "word_dom_n": 1000,
        "word_intl_n": 1000,
        "stopwords_count": 0,
        "custom_dict_count": 0,
        "dom_features": [("观光车", "8.8×", "国内高频 / 国际低频")],
        "intl_features": [("很棒", "4.3×", "国际高频 / 国内低频")],
        "network_pairs": [("熊猫", "熊猫")],
        "focus_dim": "维度情感",
        "focus_dom": 0.0,
        "focus_intl": 0.0,
        "top_dim_gaps": [("维度差异", "差距 0.00")],
        "weighted_overall_pos": 0.0,
        "weighted_overall_neu": 0.0,
        "weighted_overall_neg": 0.0,
    }

    sample_csv = Path("data/panda_sample_2000.csv")
    if sample_csv.exists():
        sample = pd.read_csv(sample_csv, encoding="utf-8-sig")
        stats["sample_n"] = len(sample)
        if "group" in sample.columns:
            stats["n_dom"] = int((sample["group"] == "国内游客").sum())
            stats["n_intl"] = int((sample["group"] == "国际游客").sum())
        if "publishDate" in sample.columns:
            dates = pd.to_datetime(sample["publishDate"], errors="coerce").dropna()
            if len(dates):
                stats["date_range"] = f"{dates.min().year} – {dates.max().year}"

    # 词频统计
    freq_csv = OUT_DIR / "1_word_freq" / "all_words_freq.csv"
    if freq_csv.exists():
        freq_df = pd.read_csv(freq_csv, encoding="utf-8-sig")
        stats["total_tokens"] = int(freq_df["count"].sum())

    sent_csv = OUT_DIR / "3_sentiment" / "sentiment_results.csv"
    if sent_csv.exists():
        sent = pd.read_csv(sent_csv, encoding="utf-8-sig")
        tr = sent["is_translated"].fillna(False).astype(bool) if "is_translated" in sent else pd.Series(False, index=sent.index)
        n_tr = int(tr.sum())
        stats["pct_translated"] = round(n_tr / max(len(sent), 1) * 100, 1)
        if "group" in sent.columns:
            intl = sent["group"] == "国际游客"
            dom = sent["group"] == "国内游客"
            stats["pct_translated_intl"] = round(float((tr & intl).sum()) / max(int(intl.sum()), 1) * 100, 1)
            stats["word_dom_n"] = int((~tr & dom).sum())
            stats["word_intl_n"] = int((~tr & intl).sum())

    weighted_csv = OUT_DIR / "3_sentiment" / "sentiment_weighted_distribution.csv"
    if weighted_csv.exists():
        weighted = pd.read_csv(weighted_csv, encoding="utf-8-sig")
        overall = weighted[weighted["分组"].astype(str).str.contains("整体")]
        if len(overall):
            r = overall.iloc[0]
            stats["weighted_overall_pos"] = float(r.get("正面占比", 0.0))
            stats["weighted_overall_neu"] = float(r.get("中性占比", 0.0))
            stats["weighted_overall_neg"] = float(r.get("负面占比", 0.0))

    for key, path in [("stopwords_count", Path("stopwords.txt")),
                      ("custom_dict_count", Path("custom_dict.txt"))]:
        if path.exists():
            stats[key] = sum(1 for line in path.read_text(encoding="utf-8").splitlines()
                             if line.strip())

    def _features(path, desc):
        if not path.exists():
            return []
        df = pd.read_csv(path, encoding="utf-8-sig").head(5)
        return [(str(r["word"]), f"{float(r['比值']):.2g}×", desc)
                for _, r in df.iterrows()]

    stats["dom_features"] = _features(
        OUT_DIR / "1_word_freq" / "distinctive_domestic.csv",
        "国内高频 / 国际低频",
    ) or stats["dom_features"]
    stats["intl_features"] = _features(
        OUT_DIR / "1_word_freq" / "distinctive_intl.csv",
        "国际高频 / 国内低频",
    ) or stats["intl_features"]

    cent_csv = OUT_DIR / "2_network" / "centrality_comparison.csv"
    if cent_csv.exists():
        cent = pd.read_csv(cent_csv, encoding="utf-8-sig").head(5)
        stats["network_pairs"] = [
            (str(r["国内词"]), str(r["国际词"])) for _, r in cent.iterrows()
        ]

    dim_csv = OUT_DIR / "4_themes" / "dimensions_table.csv"
    if dim_csv.exists():
        dim = pd.read_csv(dim_csv, encoding="utf-8-sig")
        gaps = (dim["国内情感得分"].astype(float)
                - dim["国际情感得分"].astype(float)).abs()
        dim = dim.assign(_gap=gaps).sort_values("_gap", ascending=False)
        top = dim.head(3)
        stats["top_dim_gaps"] = [
            (_clean_dim_name(r["维度名称"]), f"差距 {float(r['_gap']):.2f}")
            for _, r in top.iterrows()
        ]
        if len(top):
            r = top.iloc[0]
            stats["focus_dim"] = _clean_dim_name(r["维度名称"])
            stats["focus_dom"] = float(r["国内情感得分"])
            stats["focus_intl"] = float(r["国际情感得分"])

    return stats

# ============================================================
# 调色板 — 编辑级配色,故意压制视觉噪音
# ============================================================
INK       = RGBColor(0x1A, 0x1A, 0x1A)  # 油墨黑
SOFT      = RGBColor(0x6B, 0x6B, 0x6B)  # 中灰
SUBTLE    = RGBColor(0xA0, 0xA0, 0xA0)  # 浅灰
HAIR      = RGBColor(0xD4, 0xD4, 0xD4)  # 发丝线灰
PAPER     = RGBColor(0xFB, 0xFA, 0xF7)  # 纸白(带一点点暖)
WATERMARK = RGBColor(0xF0, 0xEE, 0xE8)  # 水印汉字色

# 仅有的两个点睛色,克制使用
ACCENT     = RGBColor(0xA0, 0x2C, 0x2A)  # 深红/朱砂(国内/警示)
ACCENT_SOFT= RGBColor(0xC6, 0xA4, 0x6E)  # 赭石/古铜(国际/平和)

WHITE = RGBColor(0xFF, 0xFF, 0xFF)

# 16:9
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# 统一安全边距(给整页留呼吸)
MARGIN_X = Inches(1.0)
MARGIN_Y_TOP = Inches(0.7)
MARGIN_Y_BOT = Inches(0.7)


# ============================================================
# 底层辅助
# ============================================================
def set_cn_font(run, name="微软雅黑", size=18, bold=False,
                color=INK, letter_spacing=None):
    """设置 run 的字体(中文也走同一字体)。"""
    rPr = run._r.get_or_add_rPr()
    rFonts = rPr.find(qn("a:rFonts"))
    if rFonts is None:
        rFonts = etree.SubElement(rPr, qn("a:rFonts"))
    rFonts.set("eastAsia", name)
    rFonts.set("ascii", name)
    rFonts.set("hAnsi", name)
    run.font.name = name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    if letter_spacing:
        # 字间距(单位 1/100 pt)
        rPr.set("spc", str(int(letter_spacing)))


def add_text(slide, left, top, width, height,
             content, *, size=16, font="微软雅黑", bold=False,
             color=INK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             letter_spacing=None, line_spacing=None):
    """添加文本框。content 可以是字符串或 list[str] 或 list[dict]。"""
    box = slide.shapes.add_textbox(left, top, width, height)
    box.fill.background()
    box.line.fill.background()
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor

    if isinstance(content, str):
        lines = [content]
    else:
        lines = content

    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if line_spacing:
            p.line_spacing = line_spacing
        if isinstance(ln, dict):
            run = p.add_run()
            set_cn_font(run,
                        name=ln.get("font", font),
                        size=ln.get("size", size),
                        bold=ln.get("bold", bold),
                        color=ln.get("color", color),
                        letter_spacing=ln.get("ls", letter_spacing))
            run.text = ln.get("text", "")
        else:
            run = p.add_run()
            set_cn_font(run, name=font, size=size, bold=bold,
                        color=color, letter_spacing=letter_spacing)
            run.text = ln
    return box


def add_richtext(slide, left, top, width, height, parts,
                 size=18, font="微软雅黑", align=PP_ALIGN.LEFT,
                 anchor=MSO_ANCHOR.TOP, line_spacing=None):
    """同一段落里混排多种样式。parts = [(text, opts), ...]"""
    box = slide.shapes.add_textbox(left, top, width, height)
    box.fill.background()
    box.line.fill.background()
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    if line_spacing:
        p.line_spacing = line_spacing
    for txt, opts in parts:
        run = p.add_run()
        set_cn_font(run,
                    name=opts.get("font", font),
                    size=opts.get("size", size),
                    bold=opts.get("bold", False),
                    color=opts.get("color", INK),
                    letter_spacing=opts.get("ls"))
        run.text = txt
    return box


def add_hairline(slide, left, top, length, *, color=HAIR, weight=0.75,
                 vertical=False):
    """加一条发丝线(默认水平)。"""
    if vertical:
        line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                          left, top, left, top + length)
    else:
        line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                          left, top, left + length, top)
    line.line.color.rgb = color
    line.line.width = Pt(weight)
    return line


def add_filled_rect(slide, left, top, width, height, fill_color,
                    line_color=None):
    """实心矩形;主要用于纸白底。"""
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 left, top, width, height)
    shp.fill.solid(); shp.fill.fore_color.rgb = fill_color
    if line_color is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line_color
        shp.line.width = Pt(0.5)
    shp.shadow.inherit = False
    return shp


def add_image(slide, image_path, left, top, width=None, height=None):
    if not Path(image_path).exists():
        add_text(slide, left, top, width or Inches(4), height or Inches(2),
                 f"[图缺失: {image_path}]", size=12, color=ACCENT)
        return None
    return slide.shapes.add_picture(str(image_path), left, top,
                                    width=width, height=height)


def blank_slide(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    # 纸白底
    add_filled_rect(s, Emu(0), Emu(0), SLIDE_W, SLIDE_H, PAPER)
    return s


# ============================================================
# 页面级公共元素
# ============================================================
def add_watermark(slide, char, left=None, top=None, size=540, color=WATERMARK):
    """大汉字水印(放在所有元素最底层)。"""
    if left is None: left = SLIDE_W - Inches(6.5)
    if top is None: top = SLIDE_H - Inches(8.2)
    add_text(slide, left, top, Inches(8), Inches(8),
             char, size=size, font="思源黑体", color=color,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)


def add_page_chrome(slide, page_no, total=9,
                    section_num=None, section_title=None):
    """页眉页脚:左下角章节编号、右下角页码、顶部细线。"""
    # 顶部 hairline(只在内页)
    add_hairline(slide, MARGIN_X, MARGIN_Y_TOP + Inches(0.7),
                 SLIDE_W - 2 * MARGIN_X)

    # 顶部章节编号 / 标题
    if section_num is not None:
        add_text(slide, MARGIN_X, MARGIN_Y_TOP, Inches(6), Inches(0.3),
                 f"{section_num:02d}",
                 size=11, color=ACCENT, bold=True, font="微软雅黑",
                 letter_spacing=250)
        if section_title:
            add_text(slide, MARGIN_X + Inches(0.6), MARGIN_Y_TOP,
                     Inches(8), Inches(0.3),
                     section_title,
                     size=11, color=SOFT, font="微软雅黑",
                     letter_spacing=200)

    # 底部:左侧学号姓名, 右侧页码
    add_text(slide, MARGIN_X, SLIDE_H - MARGIN_Y_BOT + Inches(0.05),
             Inches(8), Inches(0.3),
             f"{STUDENT_ID} · {STUDENT_NAME}",
             size=9, color=SUBTLE, font="微软雅黑", letter_spacing=150)
    add_text(slide, SLIDE_W - MARGIN_X - Inches(2),
             SLIDE_H - MARGIN_Y_BOT + Inches(0.05),
             Inches(2), Inches(0.3),
             f"{page_no:02d}  /  {total:02d}",
             size=9, color=INK, font="微软雅黑",
             align=PP_ALIGN.RIGHT, letter_spacing=200)


# ============================================================
# 第 1 页 — 封面
# ============================================================
def slide_cover(prs):
    s = blank_slide(prs)

    # 背景水印(右上一个巨大"熊"字)
    add_text(s, Inches(7.5), Inches(-1.8), Inches(8), Inches(10),
             "熊", size=620, font="思源黑体", color=WATERMARK,
             align=PP_ALIGN.LEFT)

    # 左上角课程标签
    add_text(s, MARGIN_X, Inches(0.7), Inches(8), Inches(0.3),
             "TOURISM BIG DATA / 课程实验报告",
             size=10, color=SOFT, font="微软雅黑", letter_spacing=300, bold=True)

    # 红色短线条作为视觉重音
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                             MARGIN_X, Inches(1.2),
                             Inches(0.5), Inches(0.045))
    bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()

    # 主标题 — 两行,左对齐,中文+英文小标
    add_text(s, MARGIN_X, Inches(1.6), Inches(11), Inches(1.2),
             "成都大熊猫繁育研究基地",
             size=48, bold=False, color=INK, font="微软雅黑", letter_spacing=20)
    add_text(s, MARGIN_X, Inches(2.55), Inches(11), Inches(1.0),
             "游客评论文本挖掘",
             size=48, bold=True, color=INK, font="微软雅黑", letter_spacing=20)

    # 副标题 — 用赭石与朱砂分别强调"国内"和"国际"
    add_richtext(s, MARGIN_X, Inches(3.8), Inches(11), Inches(0.5),
                 [
                     ("— ", {"size": 18, "color": SUBTLE}),
                     ("国内游客", {"size": 18, "color": ACCENT, "bold": True}),
                     (" 与 ", {"size": 18, "color": SUBTLE}),
                     ("国际游客", {"size": 18, "color": ACCENT_SOFT, "bold": True}),
                     (" 的跨文化对比研究", {"size": 18, "color": INK}),
                 ])

    # 中部数据胶囊
    cap_top = Inches(4.7)
    add_hairline(s, MARGIN_X, cap_top, Inches(11.3))
    metrics = [
        (f"{STATS.get('sample_n', 2000):,}", "条评论"),
        (f"{STATS.get('total_tokens', 0):,}", "有效词 token"),
        (f"{STATS.get('pct_translated', 0)}%", "机译识别"),
        ("Gemini 2.5", "AI 情感判别"),
    ]
    col_w = Inches(2.8)
    for i, (num, lbl) in enumerate(metrics):
        x = MARGIN_X + col_w * i
        add_text(s, x, cap_top + Inches(0.2), col_w - Inches(0.3), Inches(0.6),
                 num, size=26, bold=True, color=INK, font="微软雅黑",
                 letter_spacing=10)
        add_text(s, x, cap_top + Inches(0.95), col_w - Inches(0.3), Inches(0.4),
                 lbl, size=10, color=SOFT, font="微软雅黑", letter_spacing=200)
    add_hairline(s, MARGIN_X, cap_top + Inches(1.5), Inches(11.3))

    # 底部:学生信息(左)+ 日期(右)
    foot_y = Inches(6.45)
    add_text(s, MARGIN_X, foot_y, Inches(0.6), Inches(0.3),
             "PRESENTED BY", size=9, color=SUBTLE, font="微软雅黑",
             letter_spacing=300, bold=True)
    add_text(s, MARGIN_X, foot_y + Inches(0.3), Inches(6), Inches(0.4),
             f"{STUDENT_NAME}  ·  {STUDENT_ID}",
             size=15, color=INK, bold=True, font="微软雅黑")
    add_text(s, MARGIN_X, foot_y + Inches(0.7), Inches(6), Inches(0.3),
             f"{STUDENT_COLLEGE}  {STUDENT_MAJOR}",
             size=10, color=SOFT, font="微软雅黑", letter_spacing=100)

    add_text(s, SLIDE_W - MARGIN_X - Inches(4),
             foot_y, Inches(4), Inches(0.3),
             "DATE", size=9, color=SUBTLE,
             font="微软雅黑", letter_spacing=300, bold=True,
             align=PP_ALIGN.RIGHT)
    add_text(s, SLIDE_W - MARGIN_X - Inches(4),
             foot_y + Inches(0.3), Inches(4), Inches(0.4),
             SUBMIT_DATE,
             size=15, color=INK, bold=True, font="微软雅黑",
             align=PP_ALIGN.RIGHT)
    add_text(s, SLIDE_W - MARGIN_X - Inches(4),
             foot_y + Inches(0.7), Inches(4), Inches(0.3),
             f"指导教师  {TEACHER_NAME}",
             size=10, color=SOFT, font="微软雅黑",
             align=PP_ALIGN.RIGHT, letter_spacing=100)


# ============================================================
# 第 2 页 — 研究问题
# ============================================================
def slide_question(prs):
    s = blank_slide(prs)
    add_watermark(s, "问", left=Inches(8.2), top=Inches(-0.2), size=520)
    add_page_chrome(s, 2, section_num=1, section_title="RESEARCH QUESTION  /  研究问题")

    # 大标题左侧
    add_text(s, MARGIN_X, Inches(1.5), Inches(7), Inches(1.0),
             "三个核心问题",
             size=44, color=INK, bold=True, font="微软雅黑", letter_spacing=15)
    add_text(s, MARGIN_X, Inches(2.5), Inches(7), Inches(0.4),
             "Three guiding questions of this study",
             size=11, color=SOFT, font="微软雅黑", letter_spacing=250)

    # 三个问题块,纵向排列
    questions = [
        ("国内 vs 国际游客",
         "在评论中各自关心什么?是否存在结构性差异?"),
        ("情感维度差距",
         "哪些体验维度在两组之间情感得分差距最大?"),
        ("翻译质量影响",
         f"国际组约 {STATS['pct_translated_intl']}% 评论疑似机翻,如何识别并控制其干扰?"),
    ]
    q_top = Inches(3.6)
    line_gap = Inches(1.1)
    for i, (title, desc) in enumerate(questions):
        y = q_top + line_gap * i
        # 大编号
        add_text(s, MARGIN_X, y, Inches(0.9), Inches(0.7),
                 f"0{i+1}", size=42, color=ACCENT_SOFT,
                 bold=False, font="微软雅黑")
        # 标题 + 描述
        add_text(s, MARGIN_X + Inches(1.2), y + Inches(0.05),
                 Inches(8), Inches(0.4),
                 title, size=18, color=INK, bold=True, font="微软雅黑")
        add_text(s, MARGIN_X + Inches(1.2), y + Inches(0.5),
                 Inches(8), Inches(0.4),
                 desc, size=12, color=SOFT, font="微软雅黑")
        # 短发丝线分隔
        if i < 2:
            add_hairline(s, MARGIN_X, y + Inches(1.0), Inches(8.5))


# ============================================================
# 第 3 页 — 数据
# ============================================================
def slide_data(prs):
    s = blank_slide(prs)
    add_watermark(s, "据", left=Inches(8.0), top=Inches(-0.2), size=520)
    add_page_chrome(s, 3, section_num=2, section_title="DATA  /  数据样本")

    # Hero 数字
    add_text(s, MARGIN_X, Inches(1.5), Inches(11), Inches(0.4),
             f"携程国内版 + 携程国际版 Trip.com  · {STATS['date_range']} · 分层采样 + 总体加权",
             size=12, color=SOFT, font="微软雅黑", letter_spacing=120)

    add_text(s, MARGIN_X, Inches(2.0), Inches(8), Inches(2.5),
             "2,000",
             size=180, color=INK, bold=True, font="微软雅黑",
             letter_spacing=-20)
    add_text(s, MARGIN_X, Inches(4.4), Inches(8), Inches(0.4),
             "条 / 评论样本",
             size=18, color=INK, font="微软雅黑", letter_spacing=50)

    # 右侧:国内 / 国际 双卡(用发丝线分割,不用色块)
    right_x = Inches(9.0)
    card_w = Inches(3.5)
    card_top = Inches(2.0)
    add_hairline(s, right_x, card_top, card_w)
    add_hairline(s, right_x, card_top + Inches(2.4), card_w)
    add_hairline(s, right_x + card_w / 2, card_top, Inches(2.4), vertical=True)

    # 国内
    add_text(s, right_x + Inches(0.15), card_top + Inches(0.15),
             card_w / 2 - Inches(0.3), Inches(0.3),
             "国内游客", size=10, color=ACCENT, bold=True,
             font="微软雅黑", letter_spacing=200)
    add_text(s, right_x + Inches(0.15), card_top + Inches(0.45),
             card_w / 2 - Inches(0.3), Inches(0.9),
             "1,000", size=52, color=INK, bold=True, font="微软雅黑")
    add_text(s, right_x + Inches(0.15), card_top + Inches(1.5),
             card_w / 2 - Inches(0.3), Inches(0.4),
             "携程国内版", size=10, color=SOFT, font="微软雅黑")
    add_text(s, right_x + Inches(0.15), card_top + Inches(1.85),
             card_w / 2 - Inches(0.3), Inches(0.4),
             f"词频样本 n={STATS['word_dom_n']}",
             size=10, color=SOFT, font="微软雅黑")

    # 国际
    add_text(s, right_x + card_w/2 + Inches(0.15), card_top + Inches(0.15),
             card_w / 2 - Inches(0.3), Inches(0.3),
             "国际游客", size=10, color=ACCENT_SOFT, bold=True,
             font="微软雅黑", letter_spacing=200)
    add_text(s, right_x + card_w/2 + Inches(0.15), card_top + Inches(0.45),
             card_w / 2 - Inches(0.3), Inches(0.9),
             "1,000", size=52, color=INK, bold=True, font="微软雅黑")
    add_text(s, right_x + card_w/2 + Inches(0.15), card_top + Inches(1.5),
             card_w / 2 - Inches(0.3), Inches(0.4),
             "Trip.com", size=10, color=SOFT, font="微软雅黑")
    add_text(s, right_x + card_w/2 + Inches(0.15), card_top + Inches(1.85),
             card_w / 2 - Inches(0.3), Inches(0.4),
             f"词频样本 n={STATS['word_intl_n']}",
             size=10, color=SOFT, font="微软雅黑")

    # 底部 footnote
    add_text(s, MARGIN_X, Inches(6.0), Inches(11.3), Inches(0.4),
             f"整体口碑按清洗后 group×score 分布加权估计  ·  国际样本 {STATS['pct_translated_intl']}% 疑似机器翻译",
             size=11, color=SOFT, font="微软雅黑", letter_spacing=80)


# ============================================================
# 第 4 页 — 方法
# ============================================================
def slide_method(prs):
    s = blank_slide(prs)
    add_watermark(s, "法", left=Inches(8.5), top=Inches(-0.2), size=520)
    add_page_chrome(s, 4, section_num=3, section_title="METHODOLOGY  /  分析方法")

    add_text(s, MARGIN_X, Inches(1.5), Inches(11), Inches(0.6),
             "四阶段 Python 实现",
             size=36, color=INK, bold=True, font="微软雅黑")
    add_text(s, MARGIN_X, Inches(2.15), Inches(11), Inches(0.4),
             "等价于 ROST CM6 的核心流程,但对每个参数完全可控",
             size=12, color=SOFT, font="微软雅黑", letter_spacing=80)

    # 四个步骤横向编辑式排版
    steps = [
        ("01", "PREPROCESS", "数据预处理",
         "jieba + 自定义词典 + 停用词 + 词性过滤"),
        ("02", "FREQUENCY", "词频与词云",
         "Top50 / 差异词 / 金字塔对比图"),
        ("03", "NETWORK", "社会语义网络",
         "Top30 共现 / 度中心性 / Louvain 社区"),
        ("04", "SENTIMENT", "AI 情感与主题",
         "Gemini 篇章级 / 5 维度雷达"),
    ]
    block_top = Inches(3.4)
    col_w = Inches(2.65)
    gap = Inches(0.15)
    total = col_w * 4 + gap * 3
    left = (SLIDE_W - total) / 2
    for i, (no, en, zh, desc) in enumerate(steps):
        x = left + (col_w + gap) * i
        # 大编号
        add_text(s, x, block_top, col_w, Inches(0.9),
                 no, size=64, color=INK, bold=True, font="微软雅黑",
                 letter_spacing=-15)
        # 英文小标
        add_text(s, x, block_top + Inches(1.05), col_w, Inches(0.3),
                 en, size=10, color=ACCENT_SOFT, bold=True,
                 font="微软雅黑", letter_spacing=300)
        # 中文标
        add_text(s, x, block_top + Inches(1.4), col_w, Inches(0.4),
                 zh, size=18, color=INK, bold=True, font="微软雅黑")
        # 描述
        add_text(s, x, block_top + Inches(1.9), col_w, Inches(1.2),
                 desc, size=11, color=SOFT, font="微软雅黑",
                 line_spacing=1.4)
        # 阶段间发丝线连接(横线接续)
        if i > 0:
            add_hairline(s, x - gap, block_top + Inches(0.4),
                         gap, color=HAIR)


# ============================================================
# 第 5 页 — 发现①词频
# ============================================================
def slide_finding_freq(prs):
    s = blank_slide(prs)
    add_watermark(s, "词", left=Inches(8.5), top=Inches(-0.2), size=520)
    add_page_chrome(s, 5, section_num=4,
                    section_title="FINDING I  /  词频与特征词")

    # 大引文(pull quote)
    add_text(s, MARGIN_X, Inches(1.4), Inches(11.3), Inches(0.4),
             "FINDING 01", size=10, color=ACCENT, bold=True,
             font="微软雅黑", letter_spacing=300)

    # 用 richtext 做大字反差
    add_richtext(s, MARGIN_X, Inches(1.9), Inches(11.3), Inches(1.2),
                 [("国内重 ", {"size": 38, "color": INK, "bold": True}),
                  ("「细节」", {"size": 38, "color": ACCENT, "bold": True}),
                  (" · 国际重 ", {"size": 38, "color": INK, "bold": True}),
                  ("「整体」", {"size": 38, "color": ACCENT_SOFT, "bold": True})])

    add_text(s, MARGIN_X, Inches(3.0), Inches(11.3), Inches(0.4),
             "在 200 个候选词中,两组关注点几乎不重叠 —— 国内集中于园内运营细节,国际聚焦园区整体印象",
             size=12, color=SOFT, font="微软雅黑", letter_spacing=50)

    # 下方两列:国内特征词 / 国际特征词
    list_top = Inches(3.95)
    col1_x = MARGIN_X
    col2_x = MARGIN_X + Inches(5.7)

    # 国内
    add_text(s, col1_x, list_top, Inches(5), Inches(0.4),
             "国内特征词  TOP 5", size=10, color=ACCENT, bold=True,
             font="微软雅黑", letter_spacing=250)
    add_hairline(s, col1_x, list_top + Inches(0.4), Inches(5.2),
                 color=ACCENT, weight=1.0)
    dom_words = STATS["dom_features"]
    for i, (w, r, desc) in enumerate(dom_words):
        y = list_top + Inches(0.6) + Inches(0.45) * i
        add_text(s, col1_x, y, Inches(2.0), Inches(0.4),
                 w, size=18, color=INK, bold=True, font="微软雅黑")
        add_text(s, col1_x + Inches(2.0), y + Inches(0.05),
                 Inches(1.0), Inches(0.4),
                 r, size=12, color=ACCENT, bold=True, font="微软雅黑")
        add_text(s, col1_x + Inches(3.1), y + Inches(0.07),
                 Inches(2.1), Inches(0.4),
                 desc, size=10, color=SOFT, font="微软雅黑")

    # 国际
    add_text(s, col2_x, list_top, Inches(5), Inches(0.4),
             "国际特征词  TOP 5", size=10, color=ACCENT_SOFT, bold=True,
             font="微软雅黑", letter_spacing=250)
    add_hairline(s, col2_x, list_top + Inches(0.4), Inches(5.2),
                 color=ACCENT_SOFT, weight=1.0)
    intl_words = STATS["intl_features"]
    for i, (w, r, desc) in enumerate(intl_words):
        y = list_top + Inches(0.6) + Inches(0.45) * i
        add_text(s, col2_x, y, Inches(2.0), Inches(0.4),
                 w, size=18, color=INK, bold=True, font="微软雅黑")
        add_text(s, col2_x + Inches(2.0), y + Inches(0.05),
                 Inches(1.0), Inches(0.4),
                 r, size=12, color=ACCENT_SOFT, bold=True, font="微软雅黑")
        add_text(s, col2_x + Inches(3.1), y + Inches(0.07),
                 Inches(2.1), Inches(0.4),
                 desc, size=10, color=SOFT, font="微软雅黑")


# ============================================================
# 第 6 页 — 发现②网络
# ============================================================
def slide_finding_network(prs):
    s = blank_slide(prs)
    add_watermark(s, "网", left=Inches(8.5), top=Inches(-0.2), size=520)
    add_page_chrome(s, 6, section_num=5,
                    section_title="FINDING II  /  社会语义网络")

    add_text(s, MARGIN_X, Inches(1.4), Inches(11.3), Inches(0.4),
             "FINDING 02", size=10, color=ACCENT, bold=True,
             font="微软雅黑", letter_spacing=300)

    # Hero 标题
    add_text(s, MARGIN_X, Inches(1.9), Inches(11.3), Inches(0.8),
             "同核异构  ·  叙事分形",
             size=38, color=INK, bold=True,
             font="微软雅黑", letter_spacing=15)
    add_text(s, MARGIN_X, Inches(2.75), Inches(11.3), Inches(0.4),
             "两组网络都以「熊猫」为核心,但次级核心词截然不同 —— 国内 8 个语义社区,国际 7 个",
             size=12, color=SOFT, font="微软雅黑", letter_spacing=50)

    # 左侧:网络对比图
    add_image(s, OUT_DIR / "2_network" / "network_comparison.png",
              left=MARGIN_X, top=Inches(3.4), width=Inches(8.0))

    # 右侧:Top5 度中心性,极简表(发丝线)
    tbl_x = Inches(9.4)
    add_text(s, tbl_x, Inches(3.4), Inches(3.5), Inches(0.4),
             "度中心性  TOP 5", size=10, color=SOFT, bold=True,
             font="微软雅黑", letter_spacing=200)
    add_hairline(s, tbl_x, Inches(3.85), Inches(3.5))
    add_text(s, tbl_x, Inches(3.9), Inches(1.5), Inches(0.3),
             "国内", size=9, color=ACCENT, bold=True,
             font="微软雅黑", letter_spacing=200)
    add_text(s, tbl_x + Inches(1.75), Inches(3.9), Inches(1.5), Inches(0.3),
             "国际", size=9, color=ACCENT_SOFT, bold=True,
             font="微软雅黑", letter_spacing=200)
    add_hairline(s, tbl_x, Inches(4.2), Inches(3.5))

    pairs = STATS["network_pairs"]
    for i, (l, r) in enumerate(pairs):
        y = Inches(4.3) + Inches(0.42) * i
        add_text(s, tbl_x, y, Inches(1.5), Inches(0.4),
                 l, size=14, color=INK,
                 bold=(i == 0), font="微软雅黑")
        add_text(s, tbl_x + Inches(1.75), y, Inches(1.5), Inches(0.4),
                 r, size=14, color=INK,
                 bold=(i == 0), font="微软雅黑")
    add_hairline(s, tbl_x, Inches(4.3) + Inches(0.42) * 5, Inches(3.5))

    # 底部洞察
    add_text(s, tbl_x, Inches(6.45), Inches(3.5), Inches(0.6),
             "共有核心仅「熊猫」一词。",
             size=12, color=INK,
             font="微软雅黑", letter_spacing=30)


# ============================================================
# 第 7 页 — 发现③情感
# ============================================================
def slide_finding_sentiment(prs):
    s = blank_slide(prs)
    add_watermark(s, "感", left=Inches(8.5), top=Inches(-0.2), size=520)
    add_page_chrome(s, 7, section_num=6,
                    section_title="FINDING III  /  情感与主题")

    add_text(s, MARGIN_X, Inches(1.4), Inches(11.3), Inches(0.4),
             "FINDING 03", size=10, color=ACCENT, bold=True,
             font="微软雅黑", letter_spacing=300)

    # 大型对比数字:选取当前维度表中跨组差距最大的维度
    add_text(s, MARGIN_X, Inches(2.0), Inches(11.3), Inches(0.4),
             f"{STATS['focus_dim']}  ·  情感得分对比",
             size=12, color=SOFT, font="微软雅黑", letter_spacing=120)

    # 左侧大数字 国内
    add_text(s, MARGIN_X, Inches(2.5), Inches(5), Inches(2.0),
             _signed_score(STATS["focus_dom"]),
             size=128, color=ACCENT, bold=True, font="微软雅黑",
             letter_spacing=-15)
    add_text(s, MARGIN_X, Inches(4.55), Inches(5), Inches(0.4),
             "国内游客", size=14, color=SOFT, font="微软雅黑",
             letter_spacing=150)

    # 中间小字"vs"
    add_text(s, Inches(5.6), Inches(2.9), Inches(1), Inches(1),
             "VS", size=24, color=SUBTLE,
             font="微软雅黑", letter_spacing=200,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # 右侧大数字 国际
    add_text(s, Inches(6.5), Inches(2.5), Inches(5), Inches(2.0),
             _signed_score(STATS["focus_intl"]),
             size=128, color=ACCENT_SOFT, bold=True, font="微软雅黑",
             letter_spacing=-15, align=PP_ALIGN.RIGHT)
    add_text(s, Inches(6.5), Inches(4.55), Inches(5), Inches(0.4),
             "国际游客", size=14, color=SOFT, font="微软雅黑",
             letter_spacing=150, align=PP_ALIGN.RIGHT)

    # 分隔线
    add_hairline(s, MARGIN_X, Inches(5.2), Inches(11.3))

    # 底部三组 Top3 差距维度
    diffs = STATS["top_dim_gaps"]
    foot_top = Inches(5.5)
    col_w = Inches(3.7)
    gap = Inches(0.15)
    for i, (name, val) in enumerate(diffs):
        x = MARGIN_X + (col_w + gap) * i
        add_text(s, x, foot_top, col_w, Inches(0.3),
                 f"0{i+1}", size=9, color=ACCENT, bold=True,
                 font="微软雅黑", letter_spacing=250)
        add_text(s, x, foot_top + Inches(0.3), col_w, Inches(0.5),
                 name, size=20, color=INK, bold=True, font="微软雅黑")
        add_text(s, x, foot_top + Inches(0.85), col_w, Inches(0.4),
                 val, size=12, color=SOFT, font="微软雅黑",
                 letter_spacing=100)

    # 警示行
    add_text(s, MARGIN_X, Inches(6.6), Inches(11.3), Inches(0.3),
             f"* 整体口碑已按 group×score 加权:正面 {STATS['weighted_overall_pos']:.1f}% / 负面 {STATS['weighted_overall_neg']:.1f}%; 国际样本 {STATS['pct_translated_intl']}% 疑似机翻",
             size=10, color=SUBTLE, font="微软雅黑", letter_spacing=50)


# ============================================================
# 第 8 页 — 结论与建议
# ============================================================
def slide_conclusion(prs):
    s = blank_slide(prs)
    add_watermark(s, "论", left=Inches(8.5), top=Inches(-0.2), size=520)
    add_page_chrome(s, 8, section_num=7,
                    section_title="CONCLUSION  /  结论与建议")

    add_text(s, MARGIN_X, Inches(1.4), Inches(11), Inches(0.6),
             "三大发现 → 三条策略",
             size=36, color=INK, bold=True, font="微软雅黑", letter_spacing=15)
    add_text(s, MARGIN_X, Inches(2.05), Inches(11), Inches(0.4),
             "From cross-cultural findings to management implications",
             size=11, color=SOFT, font="微软雅黑", letter_spacing=200)

    # 横向三个结论块,极简发丝线分隔
    blocks = [
        ("01",
         "国内市场",
         "智能分流 + 增加观光车班次",
         "回应国内组在「观光车 × 排队 × 性价比」上的强烈不满"),
        ("02",
         "国际市场",
         "多语言指示 + 简化入园流程",
         "延续国际组对「公园」整体体验的高满意度"),
        ("03",
         "研究改进",
         "细分国际人群 + 机翻人工抽检",
         f"解决 {STATS['pct_translated_intl']}% 机翻文本带来的情感中性化偏差"),
    ]
    blk_top = Inches(3.0)
    col_w = (SLIDE_W - 2 * MARGIN_X - Inches(0.6)) / 3
    gap = Inches(0.3)
    for i, (no, name, action, why) in enumerate(blocks):
        x = MARGIN_X + (col_w + gap) * i
        # 大编号
        add_text(s, x, blk_top, col_w, Inches(1.0),
                 no, size=64, color=ACCENT, bold=True,
                 font="微软雅黑", letter_spacing=-10)
        # 小标签
        add_text(s, x, blk_top + Inches(1.15), col_w, Inches(0.4),
                 name, size=12, color=SOFT, bold=True,
                 font="微软雅黑", letter_spacing=200)
        # 主标
        add_text(s, x, blk_top + Inches(1.55), col_w, Inches(1.1),
                 action, size=20, color=INK, bold=True, font="微软雅黑",
                 line_spacing=1.3)
        # 描述
        add_text(s, x, blk_top + Inches(2.95), col_w, Inches(1.2),
                 why, size=11, color=SOFT, font="微软雅黑",
                 line_spacing=1.5)
        # 顶部发丝线
        add_hairline(s, x, blk_top - Inches(0.15), col_w, color=ACCENT, weight=1.2)


# ============================================================
# 第 9 页 — 致谢
# ============================================================
def slide_thanks(prs):
    s = blank_slide(prs)

    # 整页深色背景,反差大
    add_filled_rect(s, Emu(0), Emu(0), SLIDE_W, SLIDE_H, INK)

    # 大水印 "谢"
    add_text(s, Inches(5.5), Inches(-1.5), Inches(10), Inches(11),
             "谢", size=720, font="思源黑体",
             color=RGBColor(0x2A, 0x2A, 0x2A),
             align=PP_ALIGN.LEFT)

    # 顶部小字
    add_text(s, MARGIN_X, Inches(0.7), Inches(8), Inches(0.3),
             "THANK YOU  /  谢谢聆听",
             size=10, color=ACCENT_SOFT, bold=True,
             font="微软雅黑", letter_spacing=300)

    # 中部超大字 "感 谢 聆 听"
    add_text(s, MARGIN_X, Inches(2.7), Inches(11.3), Inches(1.6),
             "感谢聆听",
             size=96, color=WHITE, bold=True, font="微软雅黑",
             letter_spacing=80)

    # 红色短线
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                             MARGIN_X, Inches(4.5),
                             Inches(0.7), Inches(0.05))
    bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT_SOFT
    bar.line.fill.background()

    # 副文案
    add_text(s, MARGIN_X, Inches(4.7), Inches(11.3), Inches(0.5),
             "敬请批评指正  ·  Q & A",
             size=16, color=SUBTLE, font="微软雅黑", letter_spacing=200)

    # 底部学生信息
    add_text(s, MARGIN_X, SLIDE_H - Inches(1.8), Inches(0.6), Inches(0.3),
             "PRESENTED BY", size=9, color=SUBTLE,
             bold=True, font="微软雅黑", letter_spacing=300)
    add_text(s, MARGIN_X, SLIDE_H - Inches(1.5), Inches(8), Inches(0.4),
             f"{STUDENT_NAME}  ·  {STUDENT_ID}",
             size=14, color=WHITE, bold=True, font="微软雅黑")
    add_text(s, MARGIN_X, SLIDE_H - Inches(1.1), Inches(8), Inches(0.3),
             f"{STUDENT_COLLEGE}  {STUDENT_MAJOR}  ·  指导教师 {TEACHER_NAME}  ·  {SUBMIT_DATE}",
             size=10, color=SUBTLE, font="微软雅黑", letter_spacing=100)


# ============================================================
# 演讲备注
# ============================================================
def build_speaker_notes():
    dom_top = "、".join(w for w, _, _ in STATS["dom_features"][:4])
    intl_top = "、".join(w for w, _, _ in STATS["intl_features"][:4])
    dom_core = "、".join(l for l, _ in STATS["network_pairs"][:4])
    intl_core = "、".join(r for _, r in STATS["network_pairs"][:4])
    top_gaps = "、".join(f"{name}({gap})" for name, gap in STATS["top_dim_gaps"][:3])
    return [
        f"""各位老师好,我是旅游管理专业的罗靖坤。今天汇报的题目是《成都大熊猫繁育研究基地游客评论文本挖掘——基于国内与国际游客的跨文化对比研究》。本研究通过爬取携程国内版和国际版 Trip.com 上 {STATS['sample_n']} 条游客评论,从词频、社会网络、情感三个维度刻画国内外游客的关注差异。""",
        f"""研究问题有三个:第一,国内国际游客分别在评论里关心什么?第二,哪些体验维度在两组之间情感差距最大?第三,国际组约 {STATS['pct_translated_intl']}% 是机器翻译,这对情感判别有什么影响。""",
        f"""数据样本国内国际各 {STATS['n_dom']} 和 {STATS['n_intl']} 条,共 {STATS['sample_n']} 条,时间覆盖 {STATS['date_range']} 年。样本采用分层采样保证 1-5 星都有足够数量;整体口碑则按清洗后总体的 group×score 分布加权估计。国际样本中 {STATS['pct_translated_intl']}% 疑似机器翻译,这一点会影响后续情感判别的解读。""",
        f"""分析方法分四步:数据预处理用 jieba 加 {STATS['custom_dict_count']} 个专有名词的自定义词典和 {STATS['stopwords_count']} 个停用词;词频用 collections.Counter 并设计了"特征词比值"找两组只在一边高频的词;网络用 Top30 共现做度中心性和 Louvain 社区检测;情感用 Gemini 2.5 大模型做篇章级判别,并用总体评分分布加权校正整体比例。""",
        f"""词频对比一句话:国内重「细节」,国际重「整体」。当前国内特征词包括 {dom_top};国际特征词包括 {intl_top},显示出两组关注点和表达语境的差异。""",
        f"""语义网络的发现是"同核异构":两组都以"熊猫"为绝对核心,但次级核心词不同——国内侧前列为 {dom_core},围绕具体场景;国际侧前列为 {intl_core},围绕宏观印象。""",
        f"""情感分析最戏剧的对比发生在「{STATS['focus_dim']}」维度:国内 {_signed_score(STATS['focus_dom'])},国际 {_signed_score(STATS['focus_intl'])}。差距最大的三个维度是 {top_gaps}。需要注意的是,国际样本 {STATS['pct_translated_intl']}% 机翻,机翻会让负面表达被磨平,所以国际组的"正面"实际上可能更乐观。""",
        """综合三大发现,给出三条管理建议:国内市场重点是优化运营效率、提升性价比感知;国际市场要做好多语言无障碍服务、延续公园式整体体验;研究层面下一步是细分国际游客群体并对机翻文本做人工抽检。""",
        """感谢聆听,敬请批评指正。""",
    ]


# ============================================================
# 主流程
# ============================================================
def main():
    here = Path(__file__).parent
    os.chdir(here)
    global STATS
    STATS = _load_stats()

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    builders = [
        slide_cover,             # 1 封面
        slide_question,          # 2 研究问题
        slide_data,              # 3 数据
        slide_method,            # 4 方法
        slide_finding_freq,      # 5 发现①词频
        slide_finding_network,   # 6 发现②网络
        slide_finding_sentiment, # 7 发现③情感
        slide_conclusion,        # 8 结论
        slide_thanks,            # 9 致谢
    ]
    speaker_notes = build_speaker_notes()
    for i, build in enumerate(builders):
        build(prs)
        if i < len(speaker_notes) and speaker_notes[i]:
            prs.slides[i].notes_slide.notes_text_frame.text = speaker_notes[i]

    prs.save(str(PPT_PATH))
    size_kb = os.path.getsize(PPT_PATH) / 1024
    print(f"\n✅ PPT 已生成: {PPT_PATH}")
    print(f"   大小: {size_kb:.1f} KB · 共 {len(builders)} 页")
    print(f"   完整路径: {Path.cwd() / PPT_PATH}")


if __name__ == "__main__":
    main()
